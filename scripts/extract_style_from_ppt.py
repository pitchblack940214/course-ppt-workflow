"""Print a lightweight style summary for a course PowerPoint deck."""

from __future__ import annotations

import argparse
from collections import Counter
from pathlib import Path

from pptx import Presentation


def rgb_to_hex(rgb) -> str | None:
    if rgb is None:
        return None
    return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def summarize(path: Path, max_slides: int) -> None:
    prs = Presentation(str(path))
    fonts: Counter[str] = Counter()
    colors: Counter[str] = Counter()
    sizes: Counter[int] = Counter()

    for slide in list(prs.slides)[:max_slides]:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.name:
                        fonts[run.font.name] += 1
                    if run.font.size:
                        sizes[int(run.font.size.pt)] += 1
                    if run.font.color and run.font.color.rgb:
                        color = rgb_to_hex(run.font.color.rgb)
                        if color:
                            colors[color] += 1

    print(f"file: {path}")
    print(f"slides: {len(prs.slides)}")
    print(f"slide size: {prs.slide_width} x {prs.slide_height}")
    print("fonts:", fonts.most_common(10))
    print("font sizes:", sizes.most_common(10))
    print("font colors:", colors.most_common(10))


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--max-slides", type=int, default=8)
    args = parser.parse_args()
    summarize(args.pptx, args.max_slides)


if __name__ == "__main__":
    main()

