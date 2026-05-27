"""Validate basic PowerPoint layout safety.

Checks:
- slide count and size
- number of picture objects per slide
- objects outside the slide canvas

This cannot fully prove that text visually fits its box, but it catches the
most common layout failures before manual review.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def validate(path: Path) -> int:
    prs = Presentation(str(path))
    width, height = prs.slide_width, prs.slide_height
    failures = 0

    print(f"file: {path}")
    print(f"slides: {len(prs.slides)}")
    print(f"size: {width} x {height}")

    for index, slide in enumerate(prs.slides, 1):
        pictures = 0
        texts = 0
        out_of_bounds = []

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures += 1
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                texts += 1

            left, top, shape_width, shape_height = (
                shape.left,
                shape.top,
                shape.width,
                shape.height,
            )
            if (
                left < 0
                or top < 0
                or left + shape_width > width
                or top + shape_height > height
            ):
                out_of_bounds.append(shape.name)

        failures += len(out_of_bounds)
        print(
            f"slide {index}: texts={texts}, pictures={pictures}, "
            f"out_of_bounds={len(out_of_bounds)}"
        )
        for name in out_of_bounds:
            print(f"  - {name}")

    return failures


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    args = parser.parse_args()

    failures = validate(args.pptx)
    raise SystemExit(1 if failures else 0)


if __name__ == "__main__":
    main()

