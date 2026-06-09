"""Quality checks for generated PPTX presentation decks.

Checks:
- objects outside the slide canvas
- likely text overflow inside text boxes
- full-slide screenshots/images used as slides
- vague or hollow summary text

The text-fit and hollow-summary checks are heuristic. They are meant to catch
common problems before manual review, not to replace a visual pass.
"""

from __future__ import annotations

import argparse
import math
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


EMU_PER_INCH = 914400
POINTS_PER_INCH = 72
DEFAULT_FONT_SIZE_PT = 18.0

FULL_SLIDE_IMAGE_AREA_RATIO = 0.82
FULL_SLIDE_IMAGE_EDGE_TOLERANCE_RATIO = 0.03

VAGUE_SUMMARY_PATTERNS = [
    "综上所述",
    "总的来说",
    "整体来看",
    "具有重要意义",
    "具有启示意义",
    "进一步思考",
    "未来可以继续",
    "值得关注",
    "值得进一步研究",
    "内容总结",
    "课堂总结",
    "本节小结",
    "主要启示",
    "若干启示",
]

PLACEHOLDER_PATTERNS = [
    "这里可以输入",
    "请在此处",
    "写点小标题",
    "请输入",
    "待补充",
    "TODO",
    "TBD",
]

SUMMARY_TITLE_PATTERNS = [
    "总结",
    "小结",
    "启示",
    "讨论",
    "结论",
    "展望",
    "Summary",
    "Conclusion",
    "Discussion",
]

CONCRETE_SIGNAL_PATTERNS = [
    r"\d",
    "机制",
    "模型",
    "变量",
    "证据",
    "案例",
    "路径",
    "条件",
    "边界",
    "风险",
    "策略",
    "过程",
    "对象",
    "方法",
    "结果",
    "贡献",
    "问题",
]


@dataclass
class Issue:
    check: str
    slide: int
    message: str


def emu_to_points(value: int | float) -> float:
    return float(value) / EMU_PER_INCH * POINTS_PER_INCH


def text_width_units(text: str) -> float:
    """Approximate text width as font-size-relative units."""
    width = 0.0
    for char in text:
        if char.isspace():
            width += 0.35
        elif ord(char) < 128:
            width += 0.58
        else:
            width += 1.0
    return width


def iter_shapes(shapes) -> Iterable:
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def safe_shape_name(shape) -> str:
    return getattr(shape, "name", "<unnamed>")


def shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return shape.text.strip()


def paragraph_font_size_pt(paragraph) -> float:
    sizes = []
    for run in paragraph.runs:
        if run.font.size is not None:
            sizes.append(float(run.font.size.pt))
    return max(sizes) if sizes else DEFAULT_FONT_SIZE_PT


def estimate_text_overflow(shape) -> tuple[bool, str]:
    """Return whether a text frame likely overflows its shape."""
    if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
        return False, ""

    text_frame = shape.text_frame
    available_width = emu_to_points(
        shape.width - text_frame.margin_left - text_frame.margin_right
    )
    available_height = emu_to_points(
        shape.height - text_frame.margin_top - text_frame.margin_bottom
    )
    if available_width <= 0 or available_height <= 0:
        return True, "text frame has no usable area"

    required_height = 0.0
    widest_line = 0.0
    paragraph_count = 0

    for paragraph in text_frame.paragraphs:
        raw = "".join(run.text for run in paragraph.runs) or paragraph.text
        text = raw.strip()
        if not text:
            continue

        paragraph_count += 1
        font_size = paragraph_font_size_pt(paragraph)
        line_height = font_size * 1.25

        logical_lines = text.splitlines() or [text]
        for logical_line in logical_lines:
            estimated_width = text_width_units(logical_line) * font_size
            widest_line = max(widest_line, estimated_width)
            wrapped_lines = max(1, math.ceil(estimated_width / available_width))
            required_height += wrapped_lines * line_height

    if paragraph_count > 1:
        required_height += (paragraph_count - 1) * 3

    # PowerPoint text boxes often have a smaller nominal height than the text
    # visually occupies, especially for one-line labels. Keep this deliberately
    # conservative so the check catches likely failures instead of every tight
    # title chip.
    height_limit = max(available_height * 1.35, available_height + 10)
    height_overflow = required_height > height_limit
    line_overflow = widest_line > available_width * 3.0 and required_height > available_height * 1.15
    if height_overflow or line_overflow:
        return (
            True,
            (
                f"estimated text needs {required_height:.1f}pt high / "
                f"{widest_line:.1f}pt wide, box has {available_height:.1f}pt high / "
                f"{available_width:.1f}pt wide"
            ),
        )
    return False, ""


def object_bounds_issues(slide, slide_index: int, width: int, height: int) -> list[Issue]:
    issues = []
    for shape in iter_shapes(slide.shapes):
        if (
            shape.left < 0
            or shape.top < 0
            or shape.left + shape.width > width
            or shape.top + shape.height > height
        ):
            issues.append(
                Issue(
                    "object_bounds",
                    slide_index,
                    f"{safe_shape_name(shape)} is outside the slide canvas",
                )
            )
    return issues


def text_overflow_issues(slide, slide_index: int) -> list[Issue]:
    issues = []
    for shape in iter_shapes(slide.shapes):
        overflowing, reason = estimate_text_overflow(shape)
        if overflowing:
            issues.append(
                Issue(
                    "text_overflow",
                    slide_index,
                    f"{safe_shape_name(shape)} may overflow: {reason}",
                )
            )
    return issues


def full_slide_image_issues(slide, slide_index: int, width: int, height: int) -> list[Issue]:
    issues = []
    slide_area = width * height
    tolerance_x = width * FULL_SLIDE_IMAGE_EDGE_TOLERANCE_RATIO
    tolerance_y = height * FULL_SLIDE_IMAGE_EDGE_TOLERANCE_RATIO

    native_objects = 0
    for shape in iter_shapes(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            native_objects += 1
            continue

        area_ratio = (shape.width * shape.height) / slide_area
        near_edges = (
            shape.left <= tolerance_x
            and shape.top <= tolerance_y
            and shape.left + shape.width >= width - tolerance_x
            and shape.top + shape.height >= height - tolerance_y
        )
        if area_ratio >= FULL_SLIDE_IMAGE_AREA_RATIO or near_edges:
            issues.append(
                Issue(
                    "full_slide_image",
                    slide_index,
                    (
                        f"{safe_shape_name(shape)} covers {area_ratio:.0%} of slide; "
                        "this looks like a full-slide screenshot/image"
                    ),
                )
            )

    # A single large image plus almost no native objects is also suspicious.
    picture_count = sum(
        1 for shape in iter_shapes(slide.shapes) if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    )
    if picture_count == 1 and native_objects <= 2:
        issues.append(
            Issue(
                "full_slide_image",
                slide_index,
                "slide has one picture and very few editable objects",
            )
        )
    return issues


def has_concrete_signal(text: str) -> bool:
    return any(re.search(pattern, text, flags=re.IGNORECASE) for pattern in CONCRETE_SIGNAL_PATTERNS)


def hollow_summary_issues(slide, slide_index: int) -> list[Issue]:
    texts = [shape_text(shape) for shape in iter_shapes(slide.shapes) if shape_text(shape)]
    if not texts:
        return [
            Issue(
                "hollow_summary",
                slide_index,
                "slide has no readable text",
            )
        ]

    joined = "\n".join(texts)
    issues = []
    for pattern in PLACEHOLDER_PATTERNS:
        if pattern in joined:
            issues.append(
                Issue(
                    "hollow_summary",
                    slide_index,
                    f"placeholder text found: {pattern}",
                )
            )

    for pattern in VAGUE_SUMMARY_PATTERNS:
        if pattern in joined and not has_concrete_signal(joined):
            issues.append(
                Issue(
                    "hollow_summary",
                    slide_index,
                    f"vague summary phrase without concrete mechanism/evidence: {pattern}",
                )
            )

    title = texts[0]
    body = "\n".join(texts[1:])
    title_suggests_summary = any(pattern in title for pattern in SUMMARY_TITLE_PATTERNS)
    if title_suggests_summary and len(body) < 35:
        issues.append(
            Issue(
                "hollow_summary",
                slide_index,
                "summary/discussion slide has too little body content",
            )
        )
    if title_suggests_summary and body and not has_concrete_signal(body):
        issues.append(
            Issue(
                "hollow_summary",
                slide_index,
                "summary/discussion slide lacks concrete signals such as mechanism, evidence, method, result, or boundary",
            )
        )
    return issues


def validate(path: Path) -> int:
    prs = Presentation(str(path))
    width, height = prs.slide_width, prs.slide_height

    print(f"file: {path}")
    print(f"slides: {len(prs.slides)}")
    print(f"size: {width} x {height}")
    print(
        "checks: text_overflow, object_bounds, full_slide_image, hollow_summary"
    )

    all_issues: list[Issue] = []

    for slide_index, slide in enumerate(prs.slides, 1):
        slide_issues = []
        slide_issues.extend(object_bounds_issues(slide, slide_index, width, height))
        slide_issues.extend(text_overflow_issues(slide, slide_index))
        slide_issues.extend(full_slide_image_issues(slide, slide_index, width, height))
        slide_issues.extend(hollow_summary_issues(slide, slide_index))

        pictures = sum(
            1 for shape in iter_shapes(slide.shapes) if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        )
        text_boxes = sum(1 for shape in iter_shapes(slide.shapes) if shape_text(shape))
        print(
            f"slide {slide_index}: texts={text_boxes}, pictures={pictures}, "
            f"issues={len(slide_issues)}"
        )
        for issue in slide_issues:
            print(f"  - [{issue.check}] {issue.message}")
        all_issues.extend(slide_issues)

    print(f"total_issues: {len(all_issues)}")
    return len(all_issues)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    args = parser.parse_args()
    raise SystemExit(1 if validate(args.pptx) else 0)


if __name__ == "__main__":
    main()
