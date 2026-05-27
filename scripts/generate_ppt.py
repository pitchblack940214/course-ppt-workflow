"""Generate a small editable teaching PPT from a Markdown outline.

The script intentionally uses a public, generic style so the repository can be
shared without leaking private course decks. For real use, adapt the colors and
layout after inspecting a source course PPT.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BLUE = RGBColor(67, 115, 200)
DEEP_BLUE = RGBColor(28, 64, 125)
LIGHT_BLUE = RGBColor(238, 244, 253)
TEXT = RGBColor(35, 35, 35)
MUTED = RGBColor(92, 92, 92)


def add_textbox(slide, left, top, width, height, text, font_size=15, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    frame.text = text
    for paragraph in frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.name = "SimSun"
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = TEXT
    return box


def add_title(slide, title):
    box = slide.shapes.add_textbox(Inches(0.48), Inches(0.32), Inches(7.5), Inches(0.42))
    frame = box.text_frame
    frame.text = title
    run = frame.paragraphs[0].runs[0]
    run.font.name = "SimHei"
    run.font.bold = True
    run.font.size = Pt(23)
    run.font.color.rgb = BLUE


def add_footer(slide):
    slide.shapes.add_shape(1, Inches(1.65), Inches(6.86), Inches(10.35), Inches(0.08)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = BLUE
    slide.shapes[-1].line.color.rgb = BLUE


def add_card(slide, left, top, width, height, heading, body):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.color.rgb = RGBColor(190, 205, 230)

    heading_box = add_textbox(slide, left + Inches(0.12), top + Inches(0.08), width - Inches(0.24), Inches(0.28), heading, 14, True)
    for paragraph in heading_box.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "SimHei"
            run.font.color.rgb = DEEP_BLUE

    add_textbox(slide, left + Inches(0.12), top + Inches(0.45), width - Inches(0.24), height - Inches(0.54), body, 12)


def parse_outline(path: Path) -> list[tuple[str, list[str]]]:
    slides: list[tuple[str, list[str]]] = []
    title = None
    bullets: list[str] = []

    for raw in path.read_text(encoding="utf-8").splitlines():
        line = raw.strip()
        if line.startswith("## "):
            if title:
                slides.append((title, bullets))
            title = line[3:].strip()
            bullets = []
        elif line.startswith("- ") and title:
            bullets.append(line[2:].strip())

    if title:
        slides.append((title, bullets))
    return slides


def build(outline: Path, output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for title, bullets in parse_outline(outline):
        slide = prs.slides.add_slide(blank)
        add_title(slide, title)
        add_footer(slide)
        add_textbox(slide, Inches(0.62), Inches(1.0), Inches(3.4), Inches(0.32), "核心要点", 15, True)

        top = Inches(1.48)
        for index, bullet in enumerate(bullets[:4], 1):
            add_card(
                slide,
                Inches(0.62 + (index - 1) % 2 * 3.55),
                top + Inches(((index - 1) // 2) * 1.85),
                Inches(3.25),
                Inches(1.45),
                f"{index:02d}",
                bullet,
            )

        add_textbox(
            slide,
            Inches(8.05),
            Inches(1.05),
            Inches(4.15),
            Inches(4.6),
            "逻辑框架\n\n概念界定 -> 运行机制 -> 实务意义 -> 风险边界\n\n该页保留为可编辑文本和图形，便于教师按课程语境继续修改。",
            16,
        )

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--outline", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    args = parser.parse_args()
    build(args.outline, args.output)
    print(args.output)


if __name__ == "__main__":
    main()

